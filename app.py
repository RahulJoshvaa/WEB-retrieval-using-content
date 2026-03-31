import os
from flask import Flask, request, jsonify
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import chromadb
from chromadb.utils import embedding_functions
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = './temp_uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# --- 1. ChromaDB Setup ---
ef = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
chroma_client = chromadb.PersistentClient(path="./chroma_data")
collection = chroma_client.get_or_create_collection(name="microservice_docs", embedding_function=ef)

# --- 2. Extraction Helpers ---
def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.pdf':
            reader = PdfReader(file_path)
            return "\n".join([p.extract_text() for p in reader.pages if p.extract_text()])
        elif ext == '.docx':
            return "\n".join([para.text for para in Document(file_path).paragraphs])
        elif ext == '.pptx':
            prs = Presentation(file_path)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        return None
    except Exception as e:
        print(f"Extraction error: {e}")
        return None

def chunk_text(text, chunk_size=150, overlap=30):
    words = text.split()
    return [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size - overlap)]

# --- 3. API Endpoints ---

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    # Save temp file
    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(file_path)

    # Process
    text = extract_text(file_path)
    if not text:
        os.remove(file_path)
        return jsonify({"error": "Could not extract text or unsupported format"}), 400

    chunks = chunk_text(text)
    ids = [f"{filename}_{i}" for i in range(len(chunks))]
    metadatas = [{"file_name": filename} for _ in range(len(chunks))]

    collection.add(documents=chunks, metadatas=metadatas, ids=ids)
    
    # Cleanup temp file
    os.remove(file_path)
    
    return jsonify({
        "message": f"Successfully indexed {filename}",
        "chunks": len(chunks)
    }), 200

@app.route('/search', methods=['POST'])
def search():
    data = request.json
    query = data.get('query')
    top_k = data.get('top_k', 3)

    if not query:
        return jsonify({"error": "No query provided"}), 400

    results = collection.query(query_texts=[query], n_results=top_k)

    # Reformat for the main app to consume easily
    formatted_results = []
    for i in range(len(results['documents'][0])):
        formatted_results.append({
            "content": results['documents'][0][i],
            "file_name": results['metadatas'][0][i]['file_name'],
            "distance": results['distances'][0][i]
        })

    return jsonify({"results": formatted_results}), 200

if __name__ == '__main__':
    app.run(port=5001, debug=True)